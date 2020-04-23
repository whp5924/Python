#
import requests
from selenium import webdriver
import time
import os
from pptx import Presentation

#https://wenku.baidu.com/view/7608ad13a9956bec0975f46527d3240c8547a112.html
chrome_driver='C:\Python38-32\Scripts'
def get_html(url,return_type="text"):
    r=requests.get(url)
    r.raise_for_status()
    r.encoding=r.apparent_encoding
    if return_type=='text':
        return r.text
    else:
        return r.content
def get_picture_url(url):
    # chrome = webdriver.Chrome()
    # url = "https://baidu.com"
    # chrome.get(url)
    #
    drive=webdriver.Chrome()
    drive.get(url)
    time.sleep(5)
    img_tags=drive.find_element_by_xpath('div[@class,"ppt-image-wrap"]/img')
    img_url=[]
    for img_tag in img_tags:
        if img_tag.get_attribute('src'):
            img_urls.append(img_tag.get_attribute('src'))
        else:
            img_urls.append(img_tag.get_attribute('data-src'))
    return img_urls

def download_pictures(url):
    img_urls=get_html(url)
    print(img_urls)
    if not os.path.exists('./pictures'):
        os.makedirs('./pictures')
    path_save=[]
    for i in range(len(img_urls)):

        img=get_html(img_urls[i])
        return_type='img'
        with open('./pictures/num_%s.jpg'% str(i),'wb') as f:
            path_save.append('./pictures/num_%s.jpg' % str(i))
            f.write(img)
    return path_save

def download_word(url):
    print(url)
    global file_name
    driver=webdriver.Chrome()
    driver.get(url)
    time.sleep(5)
    #p_tages=driver.find_element_by_xpath('//div[@class,"reader-text-layer"]')
    #p_tages=driver.find_element_by_xpath('//p[contains( @class ,"reader-word")]/text()')
    #p_tages=driver.find_element_by_css_selector('//p[contains( @class ,"reader-word-layer")]/text()')
    p_tages=driver.find_elements_by_class_name('reader-word-layer')
    file_name=p_tages[2].text

    all_text=''
    for p in p_tages:
        #print(p.text)
        text=p.text
        if text=="":
            all_text=all_text+'\n'
        else:
            all_text=all_text+text
    return all_text
def generate_ppt(picture_paths):
    prs=Presentation('test.pptx')
    for picture_paths in picture_paths:
        oneSlide=prs.slides.add_slide(prs.slide_layouts[0])
        body_shapes=oneSlide.shapes.placeholders
        for index,body_shapes in enumerate(body_shapes):
            if index==0 :
                body_shape.insert_picture(picture_path)
    prs.save('baiduwenku.pptx')

def write_to_file(url,type):
    if type=='ppt':
        path_history = download_pictures(url)
        generate_ppt(path_history)
        print('下载成功......')
    elif type=='word':
        text = download_word(url)
        with open("{}.doc".format(file_name),'w',encoding='utf-8') as f:
            f.write(text)
        print('下载成功')
    else:
        print('不支付此类型文件')






if __name__=="__main__":
    print(type(__name__))
    # url1="https://wenku.baidu.com/view/7608ad13a9956bec0975f46527d3240c8547a112.html"
    url = 'https://wenku.baidu.com/view/27302dea50e79b89680203d8ce2f0066f53364e8.html'
    return_type="word"
    write_to_file(url, return_type)